# encoding: utf-8
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import requests
import time
import datetime
import os
from Tkinter import *


class App:

    def __init__(self, master):

        self.label = Label(text="Lyrics")
        self.label.grid(row=0)

        self.text = Text(master)
        self.text.grid(row=1)

        self.button = Button(master, text="Make!", command=self.make_ppt)
        self.button.grid(row=2)

    def show_message(self, text):
        top = Toplevel()
        top.title("PowerPoint Lyrics Generator")

        msg = Message(top, text=text)
        msg.pack()

        button = Button(top, text="Nice", command=top.destroy)
        button.pack()

    def make_ppt(self):
        prs = Presentation()

        textdata = self.text.get(1.0, END)

        lines = textdata.splitlines()
        linesCount = len(textdata.splitlines())

        i = 0
        while i < linesCount:

            if not lines[i]:
                i = i + 1
                continue

            slide = prs.slides.add_slide(prs.slide_layouts[6])

            url = 'https://source.unsplash.com/1600x900/?landscape'
            response = requests.get(url, stream=True)
            with open('img.png', 'wb') as f:
                f.write(response.content)
            del response

            pic = slide.shapes.add_picture('img.png', Inches(
                0), Inches(0), prs.slide_width, prs.slide_height)
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)

            txBox = slide.shapes.add_textbox(
                Inches(0), Inches(3), prs.slide_width, Inches(1.5))
            #txBox.vertical_anchor = MSO_ANCHOR.MIDDLE

            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = lines[i]
            p.alignment = PP_ALIGN.CENTER
            p.font.name = "BEBAS KAI"
            p.font.size = Pt(40)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            if i < linesCount - 1 and lines[i + 1] != "":
                p = tf.add_paragraph()
                p.text = lines[i + 1]
                p.alignment = PP_ALIGN.CENTER
                p.font.name = "BEBAS KAI"
                p.font.size = Pt(40)
                p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            txBox.fill.solid()
            txBox.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
            txBox.fill.fore_color.brightness = 0.4

            time.sleep(2.5)
            i = i + 2

        filename = lines[0].encode('utf-8') + '.pptx'
        prs.save(filename)
        os.rename(filename, lines[0] + '.pptx')
        App.show_message(self, "Done!")


root = Tk()
app = App(root)
root.resizable(width=FALSE, height=FALSE)

root.mainloop()
root.destroy()
